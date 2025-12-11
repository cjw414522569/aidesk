from docx import Document
from docx.shared import Pt, RGBColor
from openpyxl import load_workbook, Workbook
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import os

class OfficeControlMCP:
    def word_insert_text(self, filepath, text, font_size=12):
        """在Word文档中插入文本"""
        try:
            if os.path.exists(filepath):
                doc = Document(filepath)
            else:
                doc = Document()
            
            paragraph = doc.add_paragraph(text)
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
            
            doc.save(filepath)
            return f"已在Word文档中插入文本: {filepath}"
        except Exception as e:
            return f"操作失败: {str(e)}"
    
    def excel_write_cell(self, filepath, sheet_name, cell, value):
        """写入Excel单元格"""
        try:
            if os.path.exists(filepath):
                wb = load_workbook(filepath)
            else:
                wb = Workbook()
            
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                ws = wb.create_sheet(sheet_name)
            
            ws[cell] = value
            wb.save(filepath)
            return f"已写入 {sheet_name}!{cell} = {value}"
        except Exception as e:
            return f"操作失败: {str(e)}"
    
    def excel_read_cell(self, filepath, sheet_name, cell):
        """读取Excel单元格"""
        try:
            wb = load_workbook(filepath)
            ws = wb[sheet_name]
            value = ws[cell].value
            return f"{sheet_name}!{cell} = {value}"
        except Exception as e:
            return f"读取失败: {str(e)}"
    
    def ppt_add_slide(self, filepath, title, content):
        """在PPT中添加幻灯片"""
        try:
            if os.path.exists(filepath):
                prs = Presentation(filepath)
            else:
                prs = Presentation()
            
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.title.text = title
            slide.placeholders[1].text = content
            
            prs.save(filepath)
            return f"已添加幻灯片: {title}"
        except Exception as e:
            return f"操作失败: {str(e)}"
    
    def pdf_merge(self, output_file, *input_files):
        """合并PDF文件"""
        try:
            merger = PdfMerger()
            for pdf in input_files:
                merger.append(pdf)
            merger.write(output_file)
            merger.close()
            return f"已合并PDF: {output_file}"
        except Exception as e:
            return f"合并失败: {str(e)}"
    
    def pdf_split(self, input_file, output_dir, start_page, end_page):
        """拆分PDF文件"""
        try:
            reader = PdfReader(input_file)
            writer = PdfWriter()
            
            for i in range(start_page - 1, end_page):
                writer.add_page(reader.pages[i])
            
            os.makedirs(output_dir, exist_ok=True)
            output_file = os.path.join(output_dir, f"split_{start_page}-{end_page}.pdf")
            with open(output_file, 'wb') as f:
                writer.write(f)
            
            return f"已拆分PDF: {output_file}"
        except Exception as e:
            return f"拆分失败: {str(e)}"